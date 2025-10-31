import os
import json
from langchain_community.tools import ArxivQueryRun
from langchain_community.tools import DuckDuckGoSearchRun
from langchain_community.utilities import WikipediaAPIWrapper
from langchain_community.tools import WikipediaQueryRun
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from langchain_google_genai import ChatGoogleGenerativeAI
from dotenv import load_dotenv

load_dotenv()

api_key = os.getenv("GEMINI_API_KEY")
model = ChatGoogleGenerativeAI(
    model="gemini-2.5-flash", google_api_key=api_key, temperature=0.2
)
presentation = Presentation()


def generate_content_outline(topic: str, num_slides: int = 5):
    """Generate content outline using Gemini

    Args:
        topic (str): The query provided by user.
        num_slides (int): The number of slides user want in presentation.
    """

    prompt = f"""
    Create a detailed outline for a PowerPoint presentation on "{topic}" with {num_slides} slides.
    Return the response as a JSON array with the following structure:
    [
        {{
            "title": "Slide Title",
            "content": "Main content points as bullet points",
            "slide_type": "title|content|conclusion"
        }}
    ]

    Make sure the content is engaging, informative, and well-structured.
    The response must be a valid JSON array.
    """

    try:
        response = model.invoke(prompt)
        content = response.text.strip()

        # Handle potential markdown code block
        # if "```json" in content:
        #     content = content.split("```json")[1].split("```")[0].strip()
        # elif "```" in content:
        #     content = content.split("```")[1].strip()

        # Ensure we have a valid JSON string
        if not content.startswith("[") or not content.endswith("]"):
            print("Invalid JSON format received. Using fallback outline.")
            return get_fallback_outline(topic, num_slides)

        try:
            return json.loads(content)
        except json.JSONDecodeError as je:
            print(f"JSON parsing error: {je}")
            return get_fallback_outline(topic, num_slides)

    except Exception as e:
        print(f"Error generating content: {e}")
        return get_fallback_outline(topic, num_slides)


def get_fallback_outline(topic: str, num_slides: int):
    """Fallback outline if Gemini fails

    Args:
        topic (str): The query provided by user.
        num_slides (int): The number of slides user want in presentation.
    """
    return [
        {
            "title": f"Introduction to {topic}",
            "content": "• Overview of the topic\n• Key objectives\n• What to expect",
            "slide_type": "title",
        },
        {
            "title": "Background Information",
            "content": "• Historical context\n• Current state\n• Important facts",
            "slide_type": "content",
        },
        {
            "title": "Key Concepts",
            "content": "• Main principles\n• Core ideas\n• Essential knowledge",
            "slide_type": "content",
        },
        {
            "title": "Applications and Examples",
            "content": "• Real-world applications\n• Case studies\n• Practical examples",
            "slide_type": "content",
        },
        {
            "title": "Conclusion",
            "content": "• Summary of key points\n• Final thoughts\n• Next steps",
            "slide_type": "conclusion",
        },
    ]


def create_title_slide(title: str, subtitle: str = ""):
    """Create a title slide

    Args:
        title (str): Main topic on which presntation is created.
        subtitle (int): The secondry topic in the presntation.
    """

    slide_layout = presentation.slide_layouts[0]  # Title slide layout
    slide = presentation.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Set subtitle if provided
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def create_content_slide(title: str, content: str):
    """Create a content slide with bullet points

    Args:
        title (str): Topic under main presentation.
        content (str): The content of the topic.
    """

    slide_layout = presentation.slide_layouts[1]  # Title and content layout
    slide = presentation.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Set content
    content_shape = slide.placeholders[1]
    content_shape.text = content

    # Format bullet points
    text_frame = content_shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)

    return slide


def generate_presentation(
    topic: str, num_slides: int = 5, output_path: str = "generated_presentation.pptx"
):
    """Generate a complete PowerPoint presentation

    Args:
        title (str): Main topic on which presntation is created.
        subtitle (int): The secondry topic in the presntation.
    """
    print(f"Generating presentation on: {topic}")

    # Generate content outline
    outline = generate_content_outline(topic, num_slides)

    # Create slides based on outline
    for i, slide_data in enumerate(outline):
        title = slide_data["title"]
        content = slide_data["content"]
        slide_type = slide_data["slide_type"]

        print(f"Creating slide {i+1}: {title}")

        if i == 0 or slide_type == "title":
            create_title_slide(title, f"Generated by Gemini AI")
        elif slide_type == "content":
            create_content_slide(title, content)

    # Save presentation
    presentation.save(output_path)
    print(f"Presentation saved as: {output_path}")

    return output_path


def create_word_file(text: str, filename: str = "output.docx", title: str = None):
    """
    Creates a well-formatted Word file from a given string.

    Parameters:
        text (str): The content to be written to the Word document.
        filename (str): The name of the output .docx file.
        title (str, optional): Optional title to appear at the top.
    """
    # Create a new Word document
    doc = Document()

    # Set document margins
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

    # Add title if provided
    if title:
        title_paragraph = doc.add_paragraph(title)
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title_paragraph.runs[0]
        title_run.bold = True
        title_run.font.name = "Calibri"
        title_run.font.size = Pt(20)
        doc.add_paragraph()  # Add spacing after title

    # Split text into paragraphs (based on double newlines)
    paragraphs = text.split("\n\n")

    for para in paragraphs:
        if para.strip():  # Skip empty lines
            p = doc.add_paragraph(para.strip())
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            run = p.runs[0]
            run.font.name = "Calibri"
            run.font.size = Pt(12)
            p.paragraph_format.line_spacing = 1.5
            p.paragraph_format.space_after = Pt(12)

    # Save the document
    doc.save(filename)
    return f"✅ Word document '{filename}' created successfully!"


def wikipedia_search(query: str, k: int):
    """
    This tool search based on any query provided on wikipedia.

    Args:
        query (str): The query provided by user.
        k (int): The number of output you want to search on the give query.
    """
    wiki_search_tool = WikipediaQueryRun(
        api_wrapper=WikipediaAPIWrapper(top_k_results=k)
    )

    return wiki_search_tool.invoke(query)


def internet_search(query: str):
    """Run a web search"""
    search = DuckDuckGoSearchRun()
    return search.run(query)


def arxiv_search(query: str):
    """Its a search tool for research paper searching"""

    search = ArxivQueryRun()
    return search.run(query)


def create_markdown_file(filename: str, content: str):
    """
    Creates a Markdown (.md) file with the given content.

    Args:
        filename (str): The name of the Markdown file (without .md extension is okay).
        content (str): The text to write into the Markdown file.        
    """
    # Ensure the filename ends with '.md'
    if not filename.endswith(".md"):
        filename += ".md"

    try:
        with open(filename, "w", encoding="utf-8") as file:
            file.write(content)
        return f"✅ Markdown file '{filename}' created successfully!"
    except Exception as e:
        return f"❌ Error creating Markdown file: {e}"
